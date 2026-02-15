#!/usr/bin/env python3
"""Create YA design presentation as PPTX."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt
import pptx.oxml.ns as nsmap
from lxml import etree

# Colors
BG_LIGHT    = RGBColor(0xFA, 0xFA, 0xF8)
BG_DARK     = RGBColor(0x1A, 0x19, 0x15)
ACCENT      = RGBColor(0xD9, 0x77, 0x57)
TEXT_DARK   = RGBColor(0x1A, 0x19, 0x15)
TEXT_LIGHT  = RGBColor(0xFA, 0xFA, 0xF8)
MUTED       = RGBColor(0x8A, 0x89, 0x85)
BORDER      = RGBColor(0xE5, 0xE4, 0xE0)
PANEL       = RGBColor(0xF0, 0xEF, 0xEB)
DARK_PANEL  = RGBColor(0x25, 0x24, 0x20)
DARK_BORDER = RGBColor(0x35, 0x34, 0x30)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank = prs.slide_layouts[6]  # completely blank

def bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def box(slide, x, y, w, h, fill_color=None, border_color=None):
    shape = slide.shapes.add_shape(
        pptx.enum.shapes.MSO_SHAPE_TYPE.AUTO_SHAPE,
        x, y, w, h
    )
    # workaround: use add_textbox then style it
    return shape

def rect(slide, x, y, w, h, fill_color, border_color=None, radius=None):
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.oxml.ns import qn
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE = 1
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    if radius:
        # Add rounded corners via XML
        sp = shape._element
        sp_pr = sp.find(qn('p:spPr'))
        prstGeom = sp_pr.find(qn('a:prstGeom'))
        if prstGeom is not None:
            sp_pr.remove(prstGeom)
        custGeom = etree.SubElement(sp_pr, qn('a:prstGeom'))
        custGeom.set('prst', 'roundRect')
        avLst = etree.SubElement(custGeom, qn('a:avLst'))
        gd = etree.SubElement(avLst, qn('a:gd'))
        gd.set('name', 'adj')
        gd.set('fmla', f'val {radius}')
    return shape

def txt(slide, text, x, y, w, h, size=14, bold=False, color=TEXT_DARK, align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Inter"
    return tb

def dot(slide, x, y, r, color):
    from pptx.oxml.ns import qn
    shape = slide.shapes.add_shape(9, Inches(x)-Inches(r), Inches(y)-Inches(r), Inches(r*2), Inches(r*2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

# ─── SLIDE 1: Cover ───────────────────────────────────────────────────────────
s1 = prs.slides.add_slide(blank)
bg(s1, BG_LIGHT)

# Accent blob top-left
rect(s1, 0, 0, 4, 3, RGBColor(0xF5, 0xE8, 0xE0))
rect(s1, 9, 4, 5, 4, RGBColor(0xE8, 0xF0, 0xED))

# Logo circle
dot(s1, 6.665, 2.8, 0.35, ACCENT)
txt(s1, "●", 6.42, 2.45, 1, 0.6, size=28, color=ACCENT, align=PP_ALIGN.CENTER)

# Title
txt(s1, "YAgent", 3.5, 3.0, 6.3, 1.2, size=52, bold=True, color=TEXT_DARK, align=PP_ALIGN.CENTER)
txt(s1, "Персональный AI-ассистент", 3.5, 4.0, 6.3, 0.6, size=22, color=MUTED, align=PP_ALIGN.CENTER)
txt(s1, "Скачай → Запусти → Работай с AI", 3.5, 4.7, 6.3, 0.5, size=16, italic=True, color=MUTED, align=PP_ALIGN.CENTER)

# Slide number
txt(s1, "1 / 8", 12.5, 7.1, 1, 0.3, size=10, color=MUTED, align=PP_ALIGN.RIGHT)

# ─── SLIDE 2: Onboarding – Welcome ────────────────────────────────────────────
s2 = prs.slides.add_slide(blank)
bg(s2, BG_LIGHT)

# Header
txt(s2, "Шаг 1 из 2", 0, 0.3, 13.33, 0.4, size=11, color=ACCENT, align=PP_ALIGN.CENTER)

# Card
rect(s2, 3.8, 1.1, 5.7, 5.5, WHITE, BORDER, 12000)

# Inside card
txt(s2, "●", 6.1, 1.4, 1.1, 0.7, size=26, color=ACCENT, align=PP_ALIGN.CENTER)
txt(s2, "Добро пожаловать", 3.8, 2.2, 5.7, 0.8, size=24, bold=True, color=TEXT_DARK, align=PP_ALIGN.CENTER)
txt(s2, "YAgent — ваш персональный AI-ассистент.\nЗапускается локально, работает с Claude.", 3.8, 3.1, 5.7, 1.0, size=13, color=MUTED, align=PP_ALIGN.CENTER)

# Button
rect(s2, 4.8, 4.3, 3.7, 0.65, ACCENT, None, 8000)
txt(s2, "Начать →", 4.8, 4.3, 3.7, 0.65, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

txt(s2, "2 / 8", 12.5, 7.1, 1, 0.3, size=10, color=MUTED, align=PP_ALIGN.RIGHT)

# ─── SLIDE 3: Onboarding – API Key ────────────────────────────────────────────
s3 = prs.slides.add_slide(blank)
bg(s3, BG_LIGHT)

txt(s3, "Шаг 2 из 2", 0, 0.3, 13.33, 0.4, size=11, color=ACCENT, align=PP_ALIGN.CENTER)

rect(s3, 3.5, 1.0, 6.3, 5.8, WHITE, BORDER, 12000)

txt(s3, "Введите API ключ", 3.5, 1.4, 6.3, 0.7, size=22, bold=True, color=TEXT_DARK, align=PP_ALIGN.CENTER)
txt(s3, "Anthropic API Key нужен для доступа к Claude.\nПолучить: console.anthropic.com", 3.5, 2.2, 6.3, 0.8, size=12, color=MUTED, align=PP_ALIGN.CENTER)

# Input field
rect(s3, 4.0, 3.2, 5.3, 0.55, PANEL, BORDER, 5000)
txt(s3, "sk-ant-api03-••••••••••••••••••••••", 4.1, 3.22, 5.0, 0.45, size=12, color=MUTED)

txt(s3, "🔒  Ключ хранится только локально", 4.0, 3.9, 5.3, 0.4, size=11, color=MUTED, align=PP_ALIGN.CENTER)

rect(s3, 4.5, 4.6, 4.3, 0.65, ACCENT, None, 8000)
txt(s3, "Сохранить и продолжить →", 4.5, 4.6, 4.3, 0.65, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

txt(s3, "3 / 8", 12.5, 7.1, 1, 0.3, size=10, color=MUTED, align=PP_ALIGN.RIGHT)

# ─── SLIDE 4: Main UI Overview ────────────────────────────────────────────────
s4 = prs.slides.add_slide(blank)
bg(s4, BG_LIGHT)

txt(s4, "Главный интерфейс", 0.3, 0.2, 8, 0.5, size=18, bold=True, color=TEXT_DARK)
txt(s4, "Проекты (агенты) + Чаты (сессии)", 0.3, 0.65, 8, 0.35, size=12, color=MUTED)

# App shell
rect(s4, 0.3, 1.1, 12.73, 6.0, WHITE, BORDER, 8000)

# Icon rail
rect(s4, 0.3, 1.1, 0.7, 6.0, PANEL, None)
for i, icon in enumerate(["⊕", "💬", "★", "⚙"]):
    txt(s4, icon, 0.3, 1.5 + i*0.9, 0.7, 0.5, size=14, color=MUTED, align=PP_ALIGN.CENTER)

# Projects panel
rect(s4, 1.0, 1.1, 2.5, 6.0, BG_LIGHT, None)
txt(s4, "Проекты", 1.1, 1.2, 2.2, 0.4, size=11, bold=True, color=MUTED)
rect(s4, 1.05, 1.65, 2.4, 0.4, PANEL, BORDER, 5000)
txt(s4, "+ Новый проект", 1.1, 1.68, 2.3, 0.32, size=10, color=ACCENT, bold=True)
for i, proj in enumerate(["🗂 Маркетинг", "🗂 Продажи", "🗂 Аналитика"]):
    active = i == 0
    bg_c = WHITE if active else BG_LIGHT
    rect(s4, 1.05, 2.2 + i*0.7, 2.4, 0.55, bg_c, BORDER if active else None, 5000)
    txt(s4, proj, 1.15, 2.25 + i*0.7, 2.2, 0.4, size=11, color=TEXT_DARK if active else MUTED, bold=active)

# Chat area
rect(s4, 3.5, 1.1, 9.53, 6.0, WHITE, None)
txt(s4, "Маркетинг / Чат 1", 3.6, 1.2, 6, 0.4, size=13, bold=True, color=TEXT_DARK)

# Messages
rect(s4, 3.6, 1.8, 7, 0.7, PANEL, None, 8000)
txt(s4, "👤 Привет! Нужна помощь с контент-планом.", 3.7, 1.88, 6.8, 0.5, size=11, color=TEXT_DARK)
rect(s4, 3.6, 2.65, 8, 0.85, RGBColor(0xFD, 0xF0, 0xEB), None, 8000)
txt(s4, "🤖 Конечно! Вот структура контент-плана на месяц:\n1. Анализ аудитории  2. Темы постов  3. Расписание", 3.7, 2.7, 7.8, 0.75, size=11, color=TEXT_DARK)

# Input
rect(s4, 3.6, 6.3, 8.8, 0.55, PANEL, BORDER, 8000)
txt(s4, "Написать сообщение...", 3.75, 6.33, 7, 0.45, size=11, color=MUTED)
rect(s4, 12.1, 6.3, 0.55, 0.55, ACCENT, None, 8000)
txt(s4, "→", 12.1, 6.3, 0.55, 0.55, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

txt(s4, "4 / 8", 12.5, 7.1, 1, 0.3, size=10, color=MUTED, align=PP_ALIGN.RIGHT)

# ─── SLIDE 5: New Project Modal ───────────────────────────────────────────────
s5 = prs.slides.add_slide(blank)
bg(s5, BG_LIGHT)

txt(s5, "Создание нового проекта", 0.3, 0.2, 8, 0.5, size=18, bold=True, color=TEXT_DARK)

# Dimmed backdrop
rect(s5, 0.0, 0.7, 13.33, 6.8, RGBColor(0x1A, 0x19, 0x15), None)
# Lower opacity simulation — just draw overlay
from pptx.oxml.ns import qn
backdrop = s5.shapes[-1]
sp = backdrop._element
spPr = sp.find(qn('p:spPr'))
fill = spPr.find(qn('a:solidFill'))
if fill is not None:
    srgb = fill.find(qn('a:srgbClr'))
    if srgb is not None:
        alpha = etree.SubElement(srgb, qn('a:alpha'))
        alpha.set('val', '60000')

# Modal card
rect(s5, 4.2, 1.8, 4.9, 4.2, WHITE, None, 14000)
txt(s5, "Новый проект", 4.2, 2.0, 4.9, 0.55, size=18, bold=True, color=TEXT_DARK, align=PP_ALIGN.CENTER)

txt(s5, "Название", 4.5, 2.75, 1.5, 0.3, size=11, bold=True, color=MUTED)
rect(s5, 4.5, 3.05, 3.9, 0.45, PANEL, BORDER, 5000)
txt(s5, "Например: Маркетинг", 4.62, 3.08, 3.6, 0.35, size=11, color=MUTED)

txt(s5, "Описание (необязательно)", 4.5, 3.65, 3, 0.3, size=11, bold=True, color=MUTED)
rect(s5, 4.5, 3.95, 3.9, 0.45, PANEL, BORDER, 5000)
txt(s5, "Что делать в этом проекте...", 4.62, 3.98, 3.6, 0.35, size=11, color=MUTED)

rect(s5, 5.65, 4.65, 1.4, 0.45, PANEL, BORDER, 8000)
txt(s5, "Отмена", 5.65, 4.65, 1.4, 0.45, size=12, color=MUTED, align=PP_ALIGN.CENTER)
rect(s5, 7.15, 4.65, 1.6, 0.45, ACCENT, None, 8000)
txt(s5, "Создать", 7.15, 4.65, 1.6, 0.45, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

txt(s5, "5 / 8", 12.5, 7.1, 1, 0.3, size=10, color=MUTED, align=PP_ALIGN.RIGHT)

# ─── SLIDE 6: Chat in Project ─────────────────────────────────────────────────
s6 = prs.slides.add_slide(blank)
bg(s6, BG_LIGHT)

txt(s6, "Чат внутри проекта", 0.3, 0.2, 8, 0.5, size=18, bold=True, color=TEXT_DARK)
txt(s6, "Каждый чат = отдельная сессия с AI", 0.3, 0.65, 8, 0.35, size=12, color=MUTED)

rect(s6, 0.3, 1.1, 12.73, 6.0, WHITE, BORDER, 8000)
rect(s6, 0.3, 1.1, 0.7, 6.0, PANEL, None)  # rail
rect(s6, 1.0, 1.1, 2.5, 6.0, BG_LIGHT, None)  # sidebar

txt(s6, "Маркетинг", 1.1, 1.2, 2.2, 0.4, size=12, bold=True, color=TEXT_DARK)
txt(s6, "Чаты", 1.1, 1.7, 2.2, 0.35, size=10, bold=True, color=MUTED)
rect(s6, 1.05, 2.05, 2.4, 0.38, PANEL, BORDER, 5000)
txt(s6, "+ Новый чат", 1.1, 2.07, 2.3, 0.3, size=10, color=ACCENT, bold=True)
for i, chat in enumerate(["💬 Контент-план", "💬 Email-рассылка", "💬 SMM стратегия"]):
    active = i == 0
    rect(s6, 1.05, 2.55 + i*0.6, 2.4, 0.48, WHITE if active else BG_LIGHT, BORDER if active else None, 5000)
    txt(s6, chat, 1.15, 2.59 + i*0.6, 2.2, 0.35, size=10, color=TEXT_DARK if active else MUTED, bold=active)
    # Delete button on hover (shown)
    if active:
        txt(s6, "✕", 3.1, 2.62, 0.3, 0.3, size=9, color=MUTED, align=PP_ALIGN.CENTER)

# Chat header
txt(s6, "Контент-план", 3.6, 1.2, 5, 0.4, size=13, bold=True, color=TEXT_DARK)
txt(s6, "✕ Завершить чат", 10.5, 1.2, 2.3, 0.4, size=10, color=MUTED, align=PP_ALIGN.RIGHT)

# Messages
for i, (sender, msg, is_ai) in enumerate([
    ("Вы", "Составь контент-план на март для Instagram.", False),
    ("AI", "Готово! Вот план на март:\n• 1-7 марта: Знакомство с командой\n• 8-15 марта: Кейсы клиентов\n• 16-23: Советы и лайфхаки\n• 24-31: Итоги месяца", True),
]):
    bg_c = RGBColor(0xFD, 0xF0, 0xEB) if is_ai else PANEL
    rect(s6, 3.6, 1.8 + i*1.4, 8.5, 1.1, bg_c, None, 8000)
    txt(s6, f"{'🤖 AI' if is_ai else '👤 Вы'}", 3.7, 1.83 + i*1.4, 1, 0.35, size=10, bold=True, color=ACCENT if is_ai else TEXT_DARK)
    txt(s6, msg, 3.7, 2.15 + i*1.4, 8.2, 0.7, size=10, color=TEXT_DARK)

# Input
rect(s6, 3.6, 6.3, 8.8, 0.55, PANEL, BORDER, 8000)
txt(s6, "Продолжить...", 3.75, 6.33, 7, 0.45, size=11, color=MUTED)
rect(s6, 12.1, 6.3, 0.55, 0.55, ACCENT, None, 8000)
txt(s6, "→", 12.1, 6.3, 0.55, 0.55, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

txt(s6, "6 / 8", 12.5, 7.1, 1, 0.3, size=10, color=MUTED, align=PP_ALIGN.RIGHT)

# ─── SLIDE 7: Dark Theme ──────────────────────────────────────────────────────
s7 = prs.slides.add_slide(blank)
bg(s7, BG_DARK)

txt(s7, "Тёмная тема", 0.3, 0.2, 8, 0.5, size=18, bold=True, color=TEXT_LIGHT)
txt(s7, "Автоматически следует системным настройкам", 0.3, 0.65, 9, 0.35, size=12, color=RGBColor(0x60, 0x5F, 0x5A))

rect(s7, 0.3, 1.1, 12.73, 6.0, BG_DARK, DARK_BORDER, 8000)
rect(s7, 0.3, 1.1, 0.7, 6.0, DARK_PANEL, None)
rect(s7, 1.0, 1.1, 2.5, 6.0, RGBColor(0x20, 0x1F, 0x1B), None)

txt(s7, "Проекты", 1.1, 1.2, 2.2, 0.4, size=11, bold=True, color=RGBColor(0x60, 0x5F, 0x5A))
for i, proj in enumerate(["🗂 Маркетинг", "🗂 Продажи", "🗂 Аналитика"]):
    active = i == 0
    bg_c = DARK_PANEL if active else RGBColor(0x20, 0x1F, 0x1B)
    rect(s7, 1.05, 1.8 + i*0.7, 2.4, 0.55, bg_c, DARK_BORDER if active else None, 5000)
    txt(s7, proj, 1.15, 1.85 + i*0.7, 2.2, 0.4, size=11, color=TEXT_LIGHT if active else RGBColor(0x60, 0x5F, 0x5A), bold=active)

txt(s7, "Контент-план", 3.6, 1.2, 5, 0.4, size=13, bold=True, color=TEXT_LIGHT)

for i, (is_ai, msg) in enumerate([
    (False, "Составь контент-план на март."),
    (True, "Готово! Вот план на март:\n• 1-7: Знакомство с командой\n• 8-15: Кейсы клиентов"),
]):
    bg_c = RGBColor(0x2D, 0x1F, 0x15) if is_ai else DARK_PANEL
    rect(s7, 3.6, 1.8 + i*1.3, 8.5, 1.0, bg_c, None, 8000)
    txt(s7, "🤖 AI" if is_ai else "👤 Вы", 3.7, 1.83 + i*1.3, 1, 0.35, size=10, bold=True, color=ACCENT if is_ai else TEXT_LIGHT)
    txt(s7, msg, 3.7, 2.12 + i*1.3, 8.2, 0.65, size=10, color=TEXT_LIGHT)

rect(s7, 3.6, 6.3, 8.8, 0.55, DARK_PANEL, DARK_BORDER, 8000)
txt(s7, "Написать...", 3.75, 6.33, 7, 0.45, size=11, color=RGBColor(0x60, 0x5F, 0x5A))
rect(s7, 12.1, 6.3, 0.55, 0.55, ACCENT, None, 8000)
txt(s7, "→", 12.1, 6.3, 0.55, 0.55, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

txt(s7, "7 / 8", 12.5, 7.1, 1, 0.3, size=10, color=RGBColor(0x40, 0x3F, 0x3B), align=PP_ALIGN.RIGHT)

# ─── SLIDE 8: Final ───────────────────────────────────────────────────────────
s8 = prs.slides.add_slide(blank)
bg(s8, BG_LIGHT)

rect(s8, 0, 0, 4, 3, RGBColor(0xF5, 0xE8, 0xE0))
rect(s8, 9, 4, 5, 4, RGBColor(0xE8, 0xF0, 0xED))

txt(s8, "●", 6.42, 2.1, 1, 0.6, size=28, color=ACCENT, align=PP_ALIGN.CENTER)
txt(s8, "Готово к разработке", 2.5, 2.8, 8.3, 1.0, size=36, bold=True, color=TEXT_DARK, align=PP_ALIGN.CENTER)
txt(s8, "Дизайн согласован. Запускаем.", 3.5, 3.8, 6.3, 0.6, size=18, color=MUTED, align=PP_ALIGN.CENTER)

# Checklist
items = [
    ("✓", "Редизайн в стиле Anthropic", True),
    ("✓", "Управление проектами и чатами", True),
    ("✓", "Скрытое меню разработчика", True),
    ("✓", "Hard Reset скрипт", True),
    ("✓", "Тёмная тема", True),
]
for i, (check, label, done) in enumerate(items):
    color = RGBColor(0x22, 0x88, 0x55) if done else MUTED
    txt(s8, f"{check}  {label}", 4.5, 4.6 + i*0.38, 5, 0.35, size=13, color=color)

txt(s8, "8 / 8", 12.5, 7.1, 1, 0.3, size=10, color=MUTED, align=PP_ALIGN.RIGHT)

# Save
out = "/root/.openclaw/workspace/YA/design-presentation.pptx"
prs.save(out)
print(f"Saved: {out}")
